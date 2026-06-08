import io
import zipfile
from typing import List
import fitz
from fastapi.responses import StreamingResponse
from fastapi import APIRouter, UploadFile, File, Form, HTTPException
from PIL import Image
import tempfile
import os
from pdf2docx import Converter
import pdfplumber
import pandas as pd


router = APIRouter(
    prefix="/api/tools",
    tags=["PDF Tools"]
)

@router.post("/merge")
async def merge_pdfs(files: List[UploadFile] = File(...)):
    if len(files) < 2:
        raise HTTPException(status_code=400, detail="At least two PDFs are required to merge.")

    try:
        merged_pdf = fitz.open()
        for file in files:
            file_bytes = await file.read()
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            merged_pdf.insert_pdf(doc)
            doc.close()

        output = io.BytesIO()
        merged_pdf.save(output, garbage=4, deflate=True)
        merged_pdf.close()
        output.seek(0)

        return StreamingResponse(
            output,
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=RemoPDF_Merged.pdf"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Internal Processing Error: {str(e)}")


@router.post("/split")
async def split_pdf(file: UploadFile = File(...)):
    try:
        pdf_bytes = await file.read()
        src_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        zip_buffer = io.BytesIO()
        
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
            for page_idx in range(len(src_doc)):
                single_page_doc = fitz.open()
                single_page_doc.insert_pdf(src_doc, from_page=page_idx, to_page=page_idx)
                
                page_buffer = io.BytesIO()
                single_page_doc.save(page_buffer)
                single_page_doc.close()
                page_buffer.seek(0)
                
                filename = f"page_{page_idx + 1}.pdf"
                zip_file.writestr(filename, page_buffer.getvalue())
                
        src_doc.close()
        zip_buffer.seek(0)
        
        return StreamingResponse(
            zip_buffer,
            media_type="application/zip",
            headers={"Content-Disposition": f"attachment; filename=split_{file.filename}.zip"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to split PDF: {str(e)}")


@router.post("/compress")
async def compress_pdf(file: UploadFile = File(...), quality: int = Form(...)):
    try:
        file_bytes = await file.read()
        doc = fitz.open(stream=file_bytes, filetype="pdf")

        for page in doc:
            images = page.get_images()
            for img in images:
                xref = img[0]
                try:
                    pix = fitz.Pixmap(doc, xref)
                    
                    if pix.alpha:
                        pix = None
                        continue
                        
                    if pix.n >= 4: 
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    
                    img_data = pix.tobytes("jpeg", quality=quality)
                    doc.update_stream(xref, img_data)
                    doc.xref_set_key(xref, "Filter", "/DCTDecode")
                    doc.xref_set_key(xref, "ColorSpace", "/DeviceRGB" if pix.n == 3 else "/DeviceGray")
                    pix = None 
                except Exception:
                    continue

        output = io.BytesIO()
        doc.save(output, garbage=4, deflate=True)
        doc.close()
        output.seek(0)

        return StreamingResponse(
            output,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename=compressed_{file.filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Compression failed: {str(e)}")

@router.post("/image-to-pdf")
async def image_to_pdf(files: List[UploadFile] = File(...)):
    if not files:
        raise HTTPException(status_code=400, detail="At least one image file is required.")

    try:
        pdf_doc = fitz.open()
        for file in files:
            image_bytes = await file.read()
            
            # Open the image stream (handles png, jpg, webp, etc.)
            img_doc = fitz.open(stream=image_bytes, filetype="img")
            pdf_bytes = img_doc.convert_to_pdf()
            img_doc.close()
            
            # Open generated page as pdf and append to main doc
            img_pdf = fitz.open("pdf", pdf_bytes)
            pdf_doc.insert_pdf(img_pdf)
            img_pdf.close()

        output = io.BytesIO()
        pdf_doc.save(output, garbage=4, deflate=True)
        pdf_doc.close()
        output.seek(0)

        return StreamingResponse(
            output,
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=RemoPDF_Images.pdf"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Image to PDF conversion failed: {str(e)}")

@router.post("/pdf-to-images")
async def pdf_to_images(
    files: List[UploadFile] = File(...), 
    image_type: str = Form("png")
):
    if not files:
        raise HTTPException(status_code=400, detail="No PDF files provided.")
    
    # Normalize image type string format
    fmt = image_type.lower()
    if fmt not in ["png", "jpeg", "jpg"]:
        raise HTTPException(status_code=400, detail="Unsupported image format. Choose PNG or JPEG.")
    if fmt == "jpg":
        fmt = "jpeg"

    try:
        zip_buffer = io.BytesIO()
        
        # Open a ZIP archive stream to hold all page images across files
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
            for file in files:
                file_bytes = await file.read()
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                
                # Get clean base name of the file for naming output files inside the zip
                base_name = file.filename.rsplit('.', 1)[0] if '.' in file.filename else file.filename
                
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    # Render page to an image matrix (150 DPI for crisp visual balance)
                    pix = page.get_pixmap(dpi=150)
                    img_bytes = pix.tobytes(fmt)
                    
                    # Generate organized image filename inside the ZIP
                    img_filename = f"{base_name}_page_{page_num + 1}.{image_type.lower()}"
                    zip_file.writestr(img_filename, img_bytes)
                    
                doc.close()

        zip_buffer.seek(0)
        return StreamingResponse(
            zip_buffer,
            media_type="application/zip",
            headers={"Content-Disposition": "attachment; filename=RemoPDF_Images.zip"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF to Image conversion failed: {str(e)}")
        

@router.post("/get-previews")
async def get_previews(files: List[UploadFile] = File(...)):
    import base64
    try:
        result = []
        for file in files:
            file_bytes = await file.read()
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            pages_data = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                # Render to a compact 65 DPI map for fast web loading performance
                pix = page.get_pixmap(dpi=65)
                img_bytes = pix.tobytes("jpeg")
                base64_img = base64.b64encode(img_bytes).decode("utf-8")
                
                pages_data.append({
                    "page_index": page_num,
                    "thumbnail": f"data:image/jpeg;base64,{base64_img}"
                })
                
            result.append({
                "filename": file.filename,
                "pages": pages_data
            })
            doc.close()
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate previews: {str(e)}")


@router.post("/remove-pages")
async def remove_pages(
    files: List[UploadFile] = File(...),
    config: str = Form(...), 
    mode: str = Form("single")
):
    import json
    try:
        parsed_config = json.loads(config)
        processed_docs = []
        
        for idx, file in enumerate(files):
            file_bytes = await file.read()
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            
            new_doc = fitz.open()
            # Fallback to keeping all pages if file config index list mismatch occurs
            kept_indices = parsed_config[idx] if idx < len(parsed_config) else list(range(len(doc)))
            
            for p_idx in kept_indices:
                if 0 <= p_idx < len(doc):
                    new_doc.insert_pdf(doc, from_page=p_idx, to_page=p_idx)
            
            doc.close()
            processed_docs.append((file.filename, new_doc))
            
        if mode == "single":
            final_doc = fitz.open()
            for _, doc in processed_docs:
                final_doc.insert_pdf(doc)
                doc.close()
                
            output = io.BytesIO()
            final_doc.save(output, garbage=4, deflate=True)
            final_doc.close()
            output.seek(0)
            
            return StreamingResponse(
                output,
                media_type="application/pdf",
                headers={"Content-Disposition": "attachment; filename=RemoPDF_Pro_Merged.pdf"}
            )
        else:
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
                for filename, doc in processed_docs:
                    base_name = filename.rsplit('.', 1)[0] if '.' in filename else filename
                    pdf_buffer = io.BytesIO()
                    doc.save(pdf_buffer, garbage=4, deflate=True)
                    doc.close()
                    zip_file.writestr(f"{base_name}_pro.pdf", pdf_buffer.getvalue())
                    
            zip_buffer.seek(0)
            return StreamingResponse(
                zip_buffer,
                media_type="application/zip",
                headers={"Content-Disposition": "attachment; filename=RemoPDF_Pro_Package.zip"}
            )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Page manipulation failed: {str(e)}")

@router.post("/compress-images")
async def compress_images(
    files: List[UploadFile] = File(...),
    quality: int = Form(...)
):
    if not files:
        raise HTTPException(status_code=400, detail="No image files provided.")
    
    try:
        zip_buffer = io.BytesIO()
        
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
            for file in files:
                file_bytes = await file.read()
                
                # Natively open the raw memory stream using PIL (Pillow)
                img = Image.open(io.BytesIO(file_bytes))
                
                # CRITICAL FIX: Handle PNG/WEBP transparency safely
                # JPEGs can't have transparency. We place transparent images on a solid white background.
                if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
                    bg = Image.new("RGB", img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    # Use alpha channel as mask to blend cleanly
                    bg.paste(img, mask=img.split()[3])
                    img = bg
                elif img.mode != 'RGB':
                    # Catch-all for other formats
                    img = img.convert('RGB')
                
                # Compress the image natively via PIL
                img_buffer = io.BytesIO()
                img.save(img_buffer, format="JPEG", quality=quality, optimize=True)
                compressed_bytes = img_buffer.getvalue()
                
                # Format the filename cleanly
                base_name = file.filename.rsplit('.', 1)[0] if '.' in file.filename else file.filename
                img_filename = f"{base_name}_optimized.jpg"
                
                # Append to the ZIP package
                zip_file.writestr(img_filename, compressed_bytes)
                
        zip_buffer.seek(0)
        return StreamingResponse(
            zip_buffer,
            media_type="application/zip",
            headers={"Content-Disposition": "attachment; filename=RemoPDF_Premium_Compressed_Images.zip"}
        )
    except Exception as e:
        # If it fails now, it will print the exact string to your Python terminal
        print(f"CRITICAL COMPRESSION ERROR: {str(e)}") 
        raise HTTPException(status_code=500, detail=f"Image compression failed: {str(e)}")

@router.post("/pdf-to-word")
async def convert_pdf_to_word(files: List[UploadFile] = File(...)):
    if not files:
        raise HTTPException(status_code=400, detail="No PDF files provided.")
    
    try:
        zip_buffer = io.BytesIO()
        
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
            for file in files:
                file_bytes = await file.read()
                
                # Create temporary file paths
                fd_pdf, temp_pdf_path = tempfile.mkstemp(suffix=".pdf")
                fd_docx, temp_docx_path = tempfile.mkstemp(suffix=".docx")
                
                # CRITICAL WINDOWS FIX: Close raw file handles immediately to release 
                # the OS permission lock. This prevents "PermissionError" when writing.
                os.close(fd_pdf)
                os.close(fd_docx)
                
                try:
                    # Write uploaded bytes to temp PDF using standard safe context
                    with open(temp_pdf_path, 'wb') as f:
                        f.write(file_bytes)
                    
                    # Convert layout safely via physical file paths
                    cv = Converter(temp_pdf_path)
                    cv.convert(temp_docx_path, start=0, end=None) # Convert all pages
                    cv.close()
                    
                    # Read the generated docx back into memory
                    with open(temp_docx_path, "rb") as docx_file:
                        docx_bytes = docx_file.read()
                    
                    # Format filename
                    base_name = file.filename.rsplit('.', 1)[0] if '.' in file.filename else file.filename
                    docx_filename = f"{base_name}.docx"
                    
                    # Add to our download package
                    zip_file.writestr(docx_filename, docx_bytes)
                    
                finally:
                    # Always clean up physical files to prevent server storage leaks
                    if os.path.exists(temp_pdf_path):
                        os.remove(temp_pdf_path)
                    if os.path.exists(temp_docx_path):
                        os.remove(temp_docx_path)
                        
        zip_buffer.seek(0)
        return StreamingResponse(
            zip_buffer,
            media_type="application/zip",
            headers={"Content-Disposition": "attachment; filename=RemoPDF_Word_Files.zip"}
        )
    except Exception as e:
        # This will dump the exact system traceback to your Python console terminal if something else breaks
        print(f"Word Conversion Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"PDF to Word conversion failed: {str(e)}")


@router.post("/pdf-to-excel")
async def convert_pdf_to_excel(files: List[UploadFile] = File(...)):
    if not files:
        raise HTTPException(status_code=400, detail="No PDF files provided.")
    
    try:
        zip_buffer = io.BytesIO()
        
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
            for file in files:
                file_bytes = await file.read()
                
                # Secure temp files with immediate OS lock release
                fd_pdf, temp_pdf_path = tempfile.mkstemp(suffix=".pdf")
                fd_xlsx, temp_xlsx_path = tempfile.mkstemp(suffix=".xlsx")
                os.close(fd_pdf)
                os.close(fd_xlsx)
                
                try:
                    with open(temp_pdf_path, 'wb') as f:
                        f.write(file_bytes)
                    
                    # Extract tabular data safely using pdfplumber
                    all_data = []
                    with pdfplumber.open(temp_pdf_path) as pdf:
                        for page in pdf.pages:
                            # Extract tables; if none, extract standard text as fallback
                            tables = page.extract_tables()
                            if tables:
                                for table in tables:
                                    all_data.extend(table)
                                    all_data.append([]) # Visual separator row between tables
                            else:
                                text = page.extract_text()
                                if text:
                                    # Split text by newlines so it populates rows sequentially
                                    for line in text.split('\n'):
                                        all_data.append([line])
                    
                    # Convert to Excel using Pandas
                    if all_data:
                        df = pd.DataFrame(all_data)
                        df.to_excel(temp_xlsx_path, index=False, header=False)
                    else:
                        # Fallback for completely blank scanned PDFs
                        df = pd.DataFrame([["No extractable text or tables found in this document."]])
                        df.to_excel(temp_xlsx_path, index=False, header=False)
                    
                    # Read the generated Excel file back into memory
                    with open(temp_xlsx_path, "rb") as xlsx_file:
                        xlsx_bytes = xlsx_file.read()
                    
                    base_name = file.filename.rsplit('.', 1)[0] if '.' in file.filename else file.filename
                    xlsx_filename = f"{base_name}_Data.xlsx"
                    
                    zip_file.writestr(xlsx_filename, xlsx_bytes)
                    
                finally:
                    # Clean up physical files
                    if os.path.exists(temp_pdf_path):
                        os.remove(temp_pdf_path)
                    if os.path.exists(temp_xlsx_path):
                        os.remove(temp_xlsx_path)
                        
        zip_buffer.seek(0)
        return StreamingResponse(
            zip_buffer,
            media_type="application/zip",
            headers={"Content-Disposition": "attachment; filename=RemoPDF_Excel_Files.zip"}
        )
    except Exception as e:
        print(f"Excel Conversion Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"PDF to Excel conversion failed: {str(e)}")

@router.post("/add-password")
async def add_password_to_pdf(file: UploadFile = File(...), password: str = Form(...)):
    if not password:
        raise HTTPException(status_code=400, detail="Password string cannot be empty.")
        
    try:
        # Read the single uploaded PDF file into memory
        file_bytes = await file.read()
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        
        output = io.BytesIO()
        
        # CRITICAL FIX: PyMuPDF explicitly requires 'user_pw' and 'owner_pw' parameters.
        doc.save(
            output, 
            encryption=fitz.PDF_ENCRYPT_AES_256, 
            user_pw=password, 
            owner_pw=password,
            garbage=4, 
            deflate=True
        )
        doc.close()
        output.seek(0)
        
        # Deduce a secure file name
        base_name = file.filename.rsplit('.', 1)[0] if '.' in file.filename else file.filename
        protected_filename = f"{base_name}_protected.pdf"
        
        return StreamingResponse(
            output,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={protected_filename}"}
        )
    except Exception as e:
        print(f"Password Encryption Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"PDF protection failed: {str(e)}")

@router.post("/verify-password")
async def verify_password(file: UploadFile = File(...), password: str = Form(...)):
    import base64
    try:
        file_bytes = await file.read()
        doc = fitz.open(stream=file_bytes, filetype="pdf")

        if not doc.is_encrypted:
            doc.close()
            return {"is_encrypted": False, "message": "Document is not encrypted."}

        # Attempt to decrypt
        auth_success = doc.authenticate(password)
        if not auth_success:
            doc.close()
            raise HTTPException(status_code=401, detail="Incorrect password.")

        # Generate a small visual proof (thumbnail of page 1)
        page = doc[0]
        pix = page.get_pixmap(dpi=65)
        img_bytes = pix.tobytes("jpeg")
        base64_img = base64.b64encode(img_bytes).decode("utf-8")
        doc.close()

        return {"success": True, "thumbnail": f"data:image/jpeg;base64,{base64_img}"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to verify PDF: {str(e)}")

@router.post("/remove-password")
async def remove_password_from_pdf(file: UploadFile = File(...), password: str = Form(...)):
    try:
        file_bytes = await file.read()
        doc = fitz.open(stream=file_bytes, filetype="pdf")

        if doc.is_encrypted:
            auth_success = doc.authenticate(password)
            if not auth_success:
                raise HTTPException(status_code=401, detail="Incorrect password.")

        output = io.BytesIO()
        # Saving without encryption parameters strips the password entirely
        doc.save(output, garbage=4, deflate=True)
        doc.close()
        output.seek(0)

        base_name = file.filename.rsplit('.', 1)[0] if '.' in file.filename else file.filename
        unlocked_filename = f"{base_name}_unlocked.pdf"

        return StreamingResponse(
            output,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={unlocked_filename}"}
        )
    except HTTPException:
        raise
    except Exception as e:
        print(f"Decryption Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"PDF decryption failed: {str(e)}")

@router.post("/change-password")
async def change_password_of_pdf(
    file: UploadFile = File(...), 
    old_password: str = Form(...),
    new_password: str = Form(...)
):
    if not new_password:
        raise HTTPException(status_code=400, detail="New password cannot be empty.")
        
    try:
        file_bytes = await file.read()
        doc = fitz.open(stream=file_bytes, filetype="pdf")

        # Authenticate the old password if the file is currently encrypted
        if doc.is_encrypted:
            auth_success = doc.authenticate(old_password)
            if not auth_success:
                raise HTTPException(status_code=401, detail="Incorrect current password.")

        output = io.BytesIO()
        
        # Save the file with the brand new password using AES-256
        doc.save(
            output, 
            encryption=fitz.PDF_ENCRYPT_AES_256, 
            user_pw=new_password, 
            owner_pw=new_password,
            garbage=4, 
            deflate=True
        )
        doc.close()
        output.seek(0)

        base_name = file.filename.rsplit('.', 1)[0] if '.' in file.filename else file.filename
        changed_filename = f"{base_name}_updated.pdf"

        return StreamingResponse(
            output,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={changed_filename}"}
        )
    except HTTPException:
        raise
    except Exception as e:
        print(f"Password Change Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"PDF password change failed: {str(e)}")

@router.post("/pdf-to-ppt")
async def convert_pdf_to_ppt(files: List[UploadFile] = File(...)):
    from pptx import Presentation
    from pptx.util import Inches
    import io
    import zipfile
    
    if not files:
        raise HTTPException(status_code=400, detail="No PDF files provided.")
    
    try:
        zip_buffer = io.BytesIO()
        
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
            for file in files:
                file_bytes = await file.read()
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                
                # Check for password protection immediately
                if doc.is_encrypted and not doc.authenticate(""):
                    doc.close()
                    # We pass a specific ENCRYPTED tag that the React frontend will catch
                    raise HTTPException(status_code=403, detail=f"ENCRYPTED:{file.filename}")
                
                prs = Presentation()
                # Set a modern 16:9 widescreen layout
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(5.625) 
                
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    # Render page to high quality image map
                    pix = page.get_pixmap(dpi=150)
                    img_bytes = pix.tobytes("jpeg")
                    img_stream = io.BytesIO(img_bytes)
                    
                    # Add a blank layout slide (usually index 6)
                    blank_slide_layout = prs.slide_layouts[6] 
                    slide = prs.slides.add_slide(blank_slide_layout)
                    
                    # Insert the image filling the entire slide background
                    slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
                    
                doc.close()
                
                # Save presentation safely to memory buffer
                ppt_buffer = io.BytesIO()
                prs.save(ppt_buffer)
                ppt_bytes = ppt_buffer.getvalue()
                
                base_name = file.filename.rsplit('.', 1)[0] if '.' in file.filename else file.filename
                ppt_filename = f"{base_name}_Presentation.pptx"
                
                zip_file.writestr(ppt_filename, ppt_bytes)
                
        zip_buffer.seek(0)
        return StreamingResponse(
            zip_buffer,
            media_type="application/zip",
            headers={"Content-Disposition": "attachment; filename=RemoPDF_PowerPoint_Files.zip"}
        )
    except HTTPException:
        raise
    except Exception as e:
        print(f"PPT Conversion Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"PDF to PPT conversion failed: {str(e)}")